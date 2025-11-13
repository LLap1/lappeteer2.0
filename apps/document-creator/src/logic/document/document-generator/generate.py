"""
Script to replace placeholders in PowerPoint presentations with images and text
"""
from pptx import Presentation
import sys
from io import BytesIO
import struct
import json
import re
import base64


def replace_placeholders(pptx_data, parse_data):
    if isinstance(pptx_data, str):
        prs = Presentation(pptx_data)
    else:
        pptx_data.seek(0)
        prs = Presentation(pptx_data)
    
    print(f"Loaded presentation with {len(prs.slides)} slides", file=sys.stderr)
    
    string_placeholders = {}
    for string_item in parse_data.get('strings', []):
        key = string_item['key']
        value = string_item['value']
        string_placeholders[key] = value
        print(f"String placeholder: {{ {key} }} -> {value}", file=sys.stderr)
    
    map_data = parse_data.get('map')

    map_placeholders = {}
    for map_item in map_data:
        map_key = map_item['key']
        map_value = map_item['value']
        map_placeholders[map_key] = BytesIO(base64.b64decode(map_value))
    
    placeholder_pattern = re.compile(
        r'\{\{\s*([a-zA-Z_][a-zA-Z0-9_\s\.]*|[\u0590-\u05FF_][\u0590-\u05FF0-9_\s]*)\s*\}\}'
    )    
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\nProcessing slide {slide_idx}...", file=sys.stderr)
        
        shapes_to_remove = []
        images_to_add = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                original_text = shape.text
                found_placeholders = placeholder_pattern.findall(original_text)
                
                if found_placeholders:
                    print(f"  Found placeholders in shape: {found_placeholders}", file=sys.stderr)
                    
                    for placeholder_key in found_placeholders:
                        clean_key = placeholder_key.strip()
                        
                        if clean_key in map_placeholders:
                            print(f"  Replacing {{ {clean_key} }} with image", file=sys.stderr)
                            images_to_add.append({
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height,
                                'image': map_placeholders[clean_key],
                                'key': clean_key
                            })
                            shapes_to_remove.append(shape)
                            break
                    else:
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            replace_text_in_text_frame(shape.text_frame, string_placeholders, placeholder_pattern)
                        elif hasattr(shape, 'text'):
                            new_text = placeholder_pattern.sub(
                                lambda m: string_placeholders.get(m.group(1).strip(), m.group(0)),
                                shape.text
                            )
                            if new_text != original_text:
                                print(f"  Replaced text: '{original_text}' -> '{new_text}'", file=sys.stderr)
                                shape.text = new_text
            
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
            print(f"  Removed shape for image replacement", file=sys.stderr)
        
        for image_info in images_to_add:
            print(f"  Adding image for key: {image_info['key']}", file=sys.stderr)
            image_info['image'].seek(0)
            image_size_bytes = len(image_info['image'].getbuffer())
            image_size_kb = image_size_bytes / 1024
            slide.shapes.add_picture(
                image_info['image'],
                image_info['left'],
                image_info['top'],
                width=image_info['width'],
                height=image_info['height']
            )

            print(f"  Added image at position ({image_size_kb} kb - {image_info['left']}, {image_info['top']}) with size ({image_info['width']}, {image_info['height']})", file=sys.stderr)
    
    output_buffer = BytesIO()
    prs.save(output_buffer)
    output_buffer.seek(0)
    
    return output_buffer


def replace_text_in_text_frame(text_frame, placeholders, pattern):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                new_text = pattern.sub(
                    lambda m: placeholders.get(m.group(1).strip(), m.group(0)),
                    run.text
                )
                if new_text != run.text:
                    print(f"    Replaced in run: '{run.text}' -> '{new_text}'", file=sys.stderr)
                    run.text = new_text


def main():
    print("Reading template and data from stdin...", file=sys.stderr)
    
    template_size_bytes = sys.stdin.buffer.read(8)
    if len(template_size_bytes) < 8:
        print("Error: Unable to read template size", file=sys.stderr)
        sys.exit(1)
    template_size = struct.unpack('>Q', template_size_bytes)[0]
    print(f"Template size: {template_size} bytes", file=sys.stderr)
    
    template_data = BytesIO(sys.stdin.buffer.read(template_size))
    print(f"Template data read: {len(template_data.getvalue())} bytes", file=sys.stderr)
    
    parse_data_size_bytes = sys.stdin.buffer.read(8)
    if len(parse_data_size_bytes) < 8:
        print("Error: Unable to read parse data size", file=sys.stderr)
        sys.exit(1)
    parse_data_size = struct.unpack('>Q', parse_data_size_bytes)[0]
    print(f"Parse data size: {parse_data_size} bytes", file=sys.stderr)
    
    parse_data_bytes = sys.stdin.buffer.read(parse_data_size)
    print(f"Parse data read: {len(parse_data_bytes)} bytes", file=sys.stderr)
    
    try:
        parse_data = json.loads(parse_data_bytes.decode('utf-8'))
    except json.JSONDecodeError as e:
        print(f"Error: Failed to parse JSON data: {e}", file=sys.stderr)
        sys.exit(1)
    
    try:
        print(f"Parse data: {parse_data.keys()}", file=sys.stderr)
        print(f"\nReplacing placeholders...", file=sys.stderr)
        buffer = replace_placeholders(template_data, parse_data.get('data'))
        sys.stdout.buffer.write(buffer.read())
    except Exception as e:
        import traceback
        print(f"Error: {str(e)}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
