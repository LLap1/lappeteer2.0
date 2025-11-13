"""
Script to extract placeholder values from PowerPoint presentations
Finds all strings in format {{ PLACEHOLDER_NAME }}
"""
import json
import re
import sys
from io import BytesIO
import struct
from pptx import Presentation


EMU_PER_INCH = 914400
DPI = 96


def emu_to_px(value):
    return round((value / EMU_PER_INCH) * DPI, 2)


def record_placeholder(placeholders, name, width_emu, height_emu):
    if width_emu is None or height_emu is None:
        return

    width_px = emu_to_px(width_emu)
    height_px = emu_to_px(height_emu)

    if name not in placeholders:
        placeholders[name] = {
            'key': name,
            'width': width_px,
            'height': height_px
        }


def extract_placeholders_from_presentation(pptx_data):
    prs = Presentation(pptx_data)
    
    placeholders = {}
    placeholder_pattern = re.compile(r'\{\{\s*([a-zA-Z_][a-zA-Z0-9_\s]*|[\u0590-\u05FF_][\u0590-\u05FF0-9_\s]*)\s*\}\}')
    
    print(f"Analyzing presentation with {len(prs.slides)} slides", file=sys.stderr)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_idx}...", file=sys.stderr)
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            matches = placeholder_pattern.findall(run.text)
                            for match in matches:
                                match = match.strip()
                                record_placeholder(placeholders, match, getattr(shape, 'width', None), getattr(shape, 'height', None))
                                print(f"  Found placeholder: {{ {match} }} in text run", file=sys.stderr)
    
    print(f"\nTotal unique placeholders found: {len(placeholders)}", file=sys.stderr)
    
    return sorted(placeholders.values(), key=lambda item: item['key'])


def main():
    print("Reading presentation from stdin...", file=sys.stderr)
    size_bytes = sys.stdin.buffer.read(8)
    if len(size_bytes) == 8:
        size = struct.unpack('>Q', size_bytes)[0]
        print(f"Reading {size} bytes from stdin...", file=sys.stderr)
        pptx_data = BytesIO(sys.stdin.buffer.read(size))
    else:
        remaining_data = sys.stdin.buffer.read()
        pptx_data = BytesIO(size_bytes + remaining_data)
    print(f"Read {len(pptx_data.getvalue())} bytes", file=sys.stderr)
    
    try:
        placeholders = extract_placeholders_from_presentation(pptx_data)
        output_data = json.dumps(placeholders, indent=2)
        print(output_data)
    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
