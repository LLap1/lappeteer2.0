"""
Script to replace placeholders in PowerPoint presentations with images and text
"""
from pptx import Presentation
from pptx.util import Inches
import os
import argparse
import sys
from io import BytesIO
import struct
import json
import re
import base64


def replace_placeholders(pptx_data, parse_data):
    """
    Replace placeholders in the presentation with images and text.
    
    Args:
        pptx_data: BytesIO buffer or file path for the input PowerPoint file
        parse_data: Dict containing:
            - map: { type: 'map', key: string, value: string }
            - strings: [{ type: 'string', key: string, value: string }]
            - filename: string
    
    Returns:
        BytesIO buffer containing the modified presentation
    """
    # Load the presentation
    if isinstance(pptx_data, str):
        prs = Presentation(pptx_data)
    else:
        pptx_data.seek(0)
        prs = Presentation(pptx_data)
    
    print(f"Loaded presentation with {len(prs.slides)} slides", file=sys.stderr)
    
    # Build placeholder mappings
    string_placeholders = {}
    for string_item in parse_data.get('strings', []):
        key = string_item['key']
        value = string_item['value']
        string_placeholders[key] = value
        print(f"String placeholder: {{ {key} }} -> {value}", file=sys.stderr)
    
    map_data = parse_data.get('map')

    map_placeholders = {}
    for map_item in map_data:
        map_key = map_item['key']
        map_value = map_item['value']
        map_placeholders[map_key] = BytesIO(base64.b64decode(map_value))
    
    # Pattern to match {{ key }}
    placeholder_pattern = re.compile(
        r'\{\{\s*([a-zA-Z_][a-zA-Z0-9_\s\.]*|[\u0590-\u05FF_][\u0590-\u05FF0-9_\s]*)\s*\}\}'
    )    
    # Process all slides
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\nProcessing slide {slide_idx}...", file=sys.stderr)
        
        # Need to collect shapes to remove first to avoid modifying collection during iteration
        shapes_to_remove = []
        images_to_add = []
        
        for shape in slide.shapes:
            # Check if shape has text
            if hasattr(shape, 'text') and shape.text:
                original_text = shape.text
                found_placeholders = placeholder_pattern.findall(original_text)
                
                if found_placeholders:
                    print(f"  Found placeholders in shape: {found_placeholders}", file=sys.stderr)
                    
                    # Check if this is a map placeholder
                    for placeholder_key in found_placeholders:
                        clean_key = placeholder_key.strip()
                        
                        if clean_key in map_placeholders:
                            print(f"  Replacing {{ {clean_key} }} with image", file=sys.stderr)
                            # Store shape info for image replacement
                            images_to_add.append({
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height,
                                'image': map_placeholders[clean_key],
                                'key': clean_key
                            })
                            # Mark shape for removal
                            shapes_to_remove.append(shape)
                            break  # Don't process other replacements for this shape
                    else:
                        # This is a text-only shape, replace text placeholders
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            replace_text_in_text_frame(shape.text_frame, string_placeholders, placeholder_pattern)
                        elif hasattr(shape, 'text'):
                            # Simple text shape
                            new_text = placeholder_pattern.sub(
                                lambda m: string_placeholders.get(m.group(1).strip(), m.group(0)),
                                shape.text
                            )
                            if new_text != original_text:
                                print(f"  Replaced text: '{original_text}' -> '{new_text}'", file=sys.stderr)
                                shape.text = new_text
            
        
        # Remove marked shapes
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
            print(f"  Removed shape for image replacement", file=sys.stderr)
        
        # Add images
        for image_info in images_to_add:
            print(f"  Adding image for key: {image_info['key']}", file=sys.stderr)
            image_info['image'].seek(0)
            image_size_bytes = len(image_info['image'].getbuffer())
            image_size_kb = image_size_bytes / 1024
            slide.shapes.add_picture(
                image_info['image'],
                image_info['left'],
                image_info['top'],
                width=image_info['width'],
                height=image_info['height']
            )

            print(f"  Added image at position ({image_size_kb} kb - {image_info['left']}, {image_info['top']}) with size ({image_info['width']}, {image_info['height']})", file=sys.stderr)
    
    # Save to buffer
    output_buffer = BytesIO()
    prs.save(output_buffer)
    output_buffer.seek(0)
    
    return output_buffer


def replace_text_in_text_frame(text_frame, placeholders, pattern):
    """
    Replace placeholders in a text frame while preserving formatting.
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                new_text = pattern.sub(
                    lambda m: placeholders.get(m.group(1).strip(), m.group(0)),
                    run.text
                )
                if new_text != run.text:
                    print(f"    Replaced in run: '{run.text}' -> '{new_text}'", file=sys.stderr)
                    run.text = new_text


def main():
    """Main CLI entry point"""
    parser = argparse.ArgumentParser(
        description='Replace placeholders in PowerPoint presentations with images and text',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # List all images in a presentation
  python parse.py input.pptx --list

  # Replace placeholders using stdin (length-prefixed protocol)
  python parse.py --stdin > output.pptx
  
  # Replace placeholders from files
  python parse.py input.pptx -o output.pptx --data data.json
        """
    )
    
    parser.add_argument(
        'input',
        nargs='?',
        help='Input PowerPoint file path (not required with --stdin)',
        default=None
    )
    
    parser.add_argument(
        '-o', '--output',
        help='Output PowerPoint file path (optional - outputs to stdout if not provided)',
        default=None
    )
    
    parser.add_argument(
        '--data',
        help='Path to JSON file containing placeholder data',
        default=None
    )
    
    parser.add_argument(
        '--stdin',
        action='store_true',
        help='Read template and data from stdin (length-prefixed protocol)'
    )
    
    parser.add_argument(
        '-l', '--list',
        action='store_true',
        help='List all images in the presentation (no modification)'
    )
    
    args = parser.parse_args()
    
    # Validate input file exists (unless using stdin)
    if not args.stdin and args.input and not os.path.exists(args.input):
        print(f"Error: Input file '{args.input}' not found", file=sys.stderr)
        sys.exit(1)
    
    # Check for required arguments
    if not args.stdin and not args.input:
        print("Error: Input file path is required when not using --stdin", file=sys.stderr)
        sys.exit(1)

    # Get template and parse data
    if args.stdin:
        print("Reading template and data from stdin...", file=sys.stderr)
        
        # Read length-prefixed data
        # First 8 bytes: template size
        template_size_bytes = sys.stdin.buffer.read(8)
        if len(template_size_bytes) < 8:
            print("Error: Unable to read template size", file=sys.stderr)
            sys.exit(1)
        template_size = struct.unpack('>Q', template_size_bytes)[0]
        print(f"Template size: {template_size} bytes", file=sys.stderr)
        
        # Read template data
        template_data = BytesIO(sys.stdin.buffer.read(template_size))
        print(f"Template data read: {len(template_data.getvalue())} bytes", file=sys.stderr)
        
        # Next 8 bytes: parse data size
        parse_data_size_bytes = sys.stdin.buffer.read(8)
        if len(parse_data_size_bytes) < 8:
            print("Error: Unable to read parse data size", file=sys.stderr)
            sys.exit(1)
        parse_data_size = struct.unpack('>Q', parse_data_size_bytes)[0]
        print(f"Parse data size: {parse_data_size} bytes", file=sys.stderr)
        
        # Read parse data
        parse_data_bytes = sys.stdin.buffer.read(parse_data_size)
        print(f"Parse data read: {len(parse_data_bytes)} bytes", file=sys.stderr)
        
        try:
            parse_data = json.loads(parse_data_bytes.decode('utf-8'))
        except json.JSONDecodeError as e:
            print(f"Error: Failed to parse JSON data: {e}", file=sys.stderr)
            sys.exit(1)
    # Perform the replacement operation
    try:
        print(f"Parse data: {parse_data.keys()}", file=sys.stderr)
        print(f"\nReplacing placeholders...", file=sys.stderr)
        buffer = replace_placeholders(template_data, parse_data.get('data'))
        
        if args.output:
            # Save to file
            with open(args.output, 'wb') as f:
                f.write(buffer.read())
            print(f"\n✓ Success! Saved to '{args.output}'", file=sys.stderr)
        else:
            # Output to stdout as binary
            sys.stdout.buffer.write(buffer.read())
    
    except Exception as e:
        import traceback
        print(f"Error: {str(e)}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
