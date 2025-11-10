"""
Helper script to create a sample PowerPoint presentation with images for testing
"""
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont
import os


def create_sample_image(filename, width, height, color, text):
    """Create a simple colored image with text"""
    img = Image.new('RGB', (width, height), color=color)
    draw = ImageDraw.Draw(img)
    
    # Try to use a font, fall back to default if not available
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 60)
    except:
        font = ImageFont.load_default()
    
    # Get text bounding box and center it
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    
    position = ((width - text_width) // 2, (height - text_height) // 2)
    draw.text(position, text, fill='white', font=font)
    
    img.save(filename)
    print(f"Created sample image: {filename}")


def create_sample_presentation():
    """Create a sample PowerPoint presentation with images"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create sample images
    create_sample_image('original_image.png', 800, 600, 'blue', 'Original Image')
    create_sample_image('new_image.png', 800, 600, 'green', 'New Image')
    
    # Slide 1: Title slide with one image
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title = slide1.shapes.title
    
    # Add a text box for title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    textbox = slide1.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Slide 1: Single Image"
    
    # Add an image
    left = Inches(2.5)
    top = Inches(2)
    pic = slide1.shapes.add_picture('original_image.png', left, top, width=Inches(5))
    
    # Slide 2: Multiple images
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    textbox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    textbox.text_frame.text = "Slide 2: Multiple Images"
    
    # Add multiple images
    slide2.shapes.add_picture('original_image.png', Inches(0.5), Inches(2), width=Inches(4))
    slide2.shapes.add_picture('original_image.png', Inches(5.5), Inches(2), width=Inches(4))
    
    # Slide 3: No images
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    textbox.text_frame.text = "Slide 3: No Images"
    
    # Add some text content
    textbox = slide3.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
    text_frame = textbox.text_frame
    text_frame.text = "This slide has no images to replace"
    
    # Save the presentation
    prs.save('sample_presentation.pptx')
    print("\nCreated sample presentation: sample_presentation.pptx")
    print("- Slide 0: Contains 1 image")
    print("- Slide 1: Contains 2 images")
    print("- Slide 2: Contains no images")


if __name__ == "__main__":
    create_sample_presentation()
    print("\nYou can now run 'python replace_image.py' to test image replacement!")

