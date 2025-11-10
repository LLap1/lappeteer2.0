"""
Script to extract placeholder values from PowerPoint presentations
Finds all strings in format {{ PLACEHOLDER_NAME }}
"""
from pptx import Presentation
import argparse
import sys
import re
from io import BytesIO
import json
import struct


def extract_placeholders_from_presentation(pptx_data):
    """
    Extract all placeholder strings from a PowerPoint presentation.
    
    Args:
        pptx_data: BytesIO buffer or file path for the input PowerPoint file
    
    Returns:
        Set of unique placeholder names
    """
    prs = Presentation(pptx_data)
    
    placeholders = set()
    # Match either pure English or pure Hebrew placeholders (no mixing), allowing spaces
    placeholder_pattern = re.compile(r'\{\{\s*([a-zA-Z_][a-zA-Z0-9_\s]*|[\u0590-\u05FF_][\u0590-\u05FF0-9_\s]*)\s*\}\}')
    
    print(f"Analyzing presentation with {len(prs.slides)} slides", file=sys.stderr)
    
    # Iterate through all slides
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_idx}...", file=sys.stderr)
        
        # Check all shapes for text
        for shape in slide.shapes:
            # Check if shape has text
            if hasattr(shape, 'text') and shape.text:
                matches = placeholder_pattern.findall(shape.text)
                for match in matches:
                    placeholders.add(match)
                    print(f"  Found placeholder: {{ {match} }} in {shape.name}", file=sys.stderr)
            
            # Check table cells if it's a table
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            matches = placeholder_pattern.findall(cell.text)
                            for match in matches:
                                placeholders.add(match)
                                print(f"  Found placeholder: {{ {match} }} in table cell", file=sys.stderr)
            
            # Check text frame paragraphs
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            matches = placeholder_pattern.findall(run.text)
                            for match in matches:
                                placeholders.add(match)
                                print(f"  Found placeholder: {{ {match} }} in text run", file=sys.stderr)
    
    print(f"\nTotal unique placeholders found: {len(placeholders)}", file=sys.stderr)
    
    return sorted(list(placeholders))


def main():
    """Main CLI entry point"""
    parser = argparse.ArgumentParser(
        description='Extract placeholders from PowerPoint presentations',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Extract placeholders from file and output as JSON
  python extract_placeholders.py input.pptx

  # Extract placeholders from stdin
  cat input.pptx | python extract_placeholders.py --stdin

  # Save placeholders to file
  python extract_placeholders.py input.pptx -o placeholders.json
        """
    )
    
    parser.add_argument(
        'input',
        nargs='?',
        help='Input PowerPoint file path (not required with --stdin)',
        default=None
    )
    
    parser.add_argument(
        '-o', '--output',
        help='Output JSON file path (outputs to stdout if not provided)',
        default=None
    )
    
    parser.add_argument(
        '--stdin',
        action='store_true',
        help='Read PowerPoint data from stdin'
    )
    
    parser.add_argument(
        '--format',
        choices=['json', 'text', 'list'],
        default='json',
        help='Output format (default: json)'
    )
    
    args = parser.parse_args()
    
    # Check for required arguments
    if not args.stdin and not args.input:
        print("Error: Input file path is required when not using --stdin", file=sys.stderr)
        sys.exit(1)
    
    # Validate input file exists
    if not args.stdin and not os.path.exists(args.input):
        print(f"Error: Input file '{args.input}' not found", file=sys.stderr)
        sys.exit(1)
    
    # Get presentation data
    if args.stdin:
        print("Reading presentation from stdin...", file=sys.stderr)
        # Read 8 bytes for size (optional length prefix)
        size_bytes = sys.stdin.buffer.read(8)
        if len(size_bytes) == 8:
            # Has length prefix
            size = struct.unpack('>Q', size_bytes)[0]
            print(f"Reading {size} bytes from stdin...", file=sys.stderr)
            pptx_data = BytesIO(sys.stdin.buffer.read(size))
        else:
            # No length prefix, read everything
            remaining_data = sys.stdin.buffer.read()
            pptx_data = BytesIO(size_bytes + remaining_data)
        print(f"Read {len(pptx_data.getvalue())} bytes", file=sys.stderr)
    else:
        pptx_data = args.input
    
    # Extract placeholders
    try:
        placeholders = extract_placeholders_from_presentation(pptx_data)
        
        # Format output
        if args.format == 'json':
            output_data = json.dumps({
                'placeholders': placeholders,
                'count': len(placeholders)
            }, indent=2)
        elif args.format == 'list':
            output_data = '\n'.join(placeholders)
        else:  # text
            output_data = f"Found {len(placeholders)} placeholders:\n"
            for placeholder in placeholders:
                output_data += f"  - {{ {placeholder} }}\n"
        
        # Output results
        if args.output:
            with open(args.output, 'w') as f:
                f.write(output_data)
            print(f"\n✓ Success! Saved to '{args.output}'", file=sys.stderr)
        else:
            print(output_data)
    
    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()

