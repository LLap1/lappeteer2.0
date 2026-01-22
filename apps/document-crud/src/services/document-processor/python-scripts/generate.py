import sys
import json
import base64
import tempfile
import os
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

def save_data_url_to_temp(data_url: str) -> tuple[str, bool]:
    base64_data = data_url.split(',', 1)[1]
    image_data = base64.b64decode(base64_data)
    temp_path = tempfile.mktemp(suffix='.png')
    with open(temp_path, 'wb') as f:
        f.write(image_data)
    return temp_path, True

def get_image_path(source: str) -> tuple[str, bool]:
    if source.startswith('data:'):
        return save_data_url_to_temp(source)
    if not os.path.exists(source):
        raise FileNotFoundError(f"Image file not found: {source}")
    return source, False

class ImageLayer:
    def __init__(self, path: str, offset_x: float, offset_y: float, width: float, height: float):
        self.path = path
        self.offset_x = offset_x
        self.offset_y = offset_y
        self.width = width
        self.height = height

class TextValue:
    def __init__(self, text: str, fit_mode: str = 'shrink'):
        self.text = text
        self.fit_mode = fit_mode

def prepare_data(placeholder_data: list[dict]) -> tuple[dict[str, TextValue], dict[str, tuple[list[ImageLayer], int | None]], list[str]]:
    text_values: dict[str, TextValue] = {}
    image_values: dict[str, tuple[list[ImageLayer], int | None]] = {}
    temp_files_to_cleanup: list[str] = []

    for item in placeholder_data:
        key = item['key'].strip()
        placeholder_type = item['type']
        value = item['value']
        rotation = item.get('rotation')

        if placeholder_type == 'text':
            text_data = json.loads(value)
            text_values[key] = TextValue(
                text_data.get('text', ''),
                text_data.get('fitMode', 'shrink'),
            )
        elif placeholder_type in ('image', 'map'):
            layers_data = json.loads(value)
            layers = []
            for layer_data in layers_data:
                if isinstance(layer_data, str):
                    file_path, should_cleanup = get_image_path(layer_data)
                    layers.append(ImageLayer(file_path, 0, 0, 0, 0))
                else:
                    file_path, should_cleanup = get_image_path(layer_data['path'])
                    layers.append(ImageLayer(
                        file_path,
                        layer_data.get('offsetX', 0),
                        layer_data.get('offsetY', 0),
                        layer_data.get('width', 0),
                        layer_data.get('height', 0),
                    ))
                if should_cleanup:
                    temp_files_to_cleanup.append(file_path)
            
            image_values[key] = (layers, rotation)

    return text_values, image_values, temp_files_to_cleanup

def get_run_color(run) -> RGBColor | None:
    try:
        if run.font.color.rgb:
            return run.font.color.rgb
    except:
        pass
    return None

def get_type_from_color(color: RGBColor | None) -> str | None:
    if color is None:
        return None
    
    r, g, b = color[0], color[1], color[2]
    
    if r > 150 and g < 100 and b < 100:
        return 'text'
    if r < 100 and g < 100 and b > 150:
        return 'image'
    if r < 100 and g > 150 and b < 100:
        return 'map'
    
    return None

def is_placeholder_run(run) -> bool:
    font = run.font
    return font.bold and font.italic and font.underline

def get_font_name(run, paragraph) -> str:
    if run.font.name:
        return run.font.name
    if paragraph.font.name:
        return paragraph.font.name
    return 'Arial'

def get_original_font_size(run) -> int:
    if run.font.size:
        return int(run.font.size.pt)
    return 18

def estimate_text_needs_wrap(text: str, font_size_pt: float, width_pt: float) -> int:
    avg_char_width = font_size_pt * 0.5
    chars_per_line = max(1, int(width_pt / avg_char_width))
    num_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
    return num_lines

def estimate_text_height(num_lines: int, font_size_pt: float) -> float:
    line_height = font_size_pt * 1.1
    return num_lines * line_height

def calculate_best_font_size(text: str, original_size: int, width_pt: float, height_pt: float, min_size: int = 6) -> int:
    if width_pt <= 0 or height_pt <= 0:
        return original_size
    
    for size in range(original_size, min_size - 1, -1):
        num_lines = estimate_text_needs_wrap(text, size, width_pt)
        text_height = estimate_text_height(num_lines, size)
        
        if text_height <= height_pt:
            return size
    
    return min_size

def process_text_frame(text_frame, text_values: dict[str, TextValue], is_cell: bool = False, cell_width: float = 0, cell_height: float = 0) -> None:
    font_name = 'Arial'
    original_size = 18
    modified = False
    new_text = ""

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not is_placeholder_run(run):
                continue
            
            key = run.text.strip()
            if key not in text_values:
                continue
            
            color = get_run_color(run)
            placeholder_type = get_type_from_color(color)
            
            if placeholder_type != 'text':
                continue
            
            font_name = get_font_name(run, paragraph)
            original_size = get_original_font_size(run)
            text_value = text_values[key]
            new_text = text_value.text
            fit_mode = text_value.fit_mode
            run.text = new_text
            run.font.bold = False
            run.font.italic = False
            run.font.underline = False
            run.font.color.rgb = RGBColor(0, 0, 0)
   


def find_image_placeholder_in_shape(shape, image_values: dict[str, tuple[list[ImageLayer], int | None]]) -> tuple[str, str] | None:
    if not hasattr(shape, 'text_frame'):
        return None
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not is_placeholder_run(run):
                continue
            
            key = run.text.strip()
            if key not in image_values:
                continue
            
            color = get_run_color(run)
            placeholder_type = get_type_from_color(color)
            
            if placeholder_type in ('image', 'map'):
                return key, placeholder_type
    
    return None

def rotate_image(image_path: str, rotation: int) -> str:
    img = Image.open(image_path).convert('RGBA')
    rotated = img.rotate(-rotation, expand=True, resample=Image.Resampling.BICUBIC, fillcolor=(0, 0, 0, 0))
    temp_path = tempfile.mktemp(suffix='.png')
    rotated.save(temp_path, 'PNG')
    return temp_path

def get_shape_z_index(shape) -> int:
    parent = shape.element.getparent()
    return list(parent).index(shape.element)

def move_shape_to_z_index(shape, z_index: int) -> None:
    parent = shape.element.getparent()
    parent.remove(shape.element)
    parent.insert(z_index, shape.element)

def process_shape(shape, slide, text_values: dict[str, TextValue], image_values: dict[str, tuple[list[ImageLayer], int | None]]) -> bool:
    placeholder_info = find_image_placeholder_in_shape(shape, image_values)
    
    if placeholder_info:
        key, _ = placeholder_info
        layers, rotation = image_values[key]
        placeholder_left = shape.left
        placeholder_top = shape.top
        placeholder_width = shape.width
        placeholder_height = shape.height

        z_index = get_shape_z_index(shape)
        shape.element.getparent().remove(shape.element)
        
        base_layer = layers[0] if layers else None
        if not base_layer or base_layer.width == 0:
            map_width = placeholder_width
            map_height = placeholder_height
            scale_x = 1.0
            scale_y = 1.0
        else:
            map_width = base_layer.width
            map_height = base_layer.height
            scale_x = placeholder_width / map_width
            scale_y = placeholder_height / map_height

        for i, layer in enumerate(layers):
            if layer.width == 0 or layer.height == 0:
                layer_left = placeholder_left
                layer_top = placeholder_top
                layer_width = placeholder_width
                layer_height = placeholder_height
            else:
                layer_left = placeholder_left + int(layer.offset_x * scale_x)
                layer_top = placeholder_top + int(layer.offset_y * scale_y)
                layer_width = int(layer.width * scale_x)
                layer_height = int(layer.height * scale_y)

            if rotation:
                rotated_path = rotate_image(layer.path, rotation)
                picture = slide.shapes.add_picture(rotated_path, layer_left, layer_top, layer_width, layer_height)
                os.unlink(rotated_path)
            else:
                picture = slide.shapes.add_picture(layer.path, layer_left, layer_top, layer_width, layer_height)
            
            move_shape_to_z_index(picture, z_index + i)

        return True

    if hasattr(shape, 'text_frame'):
        process_text_frame(shape.text_frame, text_values)
    
    return False


def process_table(table, text_values: dict[str, TextValue], image_values: dict[str, tuple[list[ImageLayer], int | None]], slide) -> None:
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            if hasattr(cell, 'text_frame'):
                cell_width = 0
                cell_height = 0
                
                try:
                    if table.columns[col_idx].width:
                        cell_width = table.columns[col_idx].width.pt - 12
                    if row.height:
                        cell_height = row.height.pt - 12
                except Exception:
                    pass
                
                process_text_frame(cell.text_frame, text_values, is_cell=True, cell_width=cell_width, cell_height=cell_height)

def process_all_shapes(shapes, slide, text_values: dict[str, TextValue], image_values: dict[str, tuple[list[ImageLayer], int | None]]) -> None:
    shapes_list = list(shapes)
    
    for shape in shapes_list:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            if hasattr(shape, 'shapes'):
                process_all_shapes(shape.shapes, slide, text_values, image_values)
        elif hasattr(shape, 'has_table') and shape.has_table:
            process_table(shape.table, text_values, image_values, slide)
        else:
            process_shape(shape, slide, text_values, image_values)

def cleanup_temp_files(temp_files: list[str]) -> None:
    for file_path in temp_files:
        try:
            os.unlink(file_path)
        except:
            pass

def generate(file_path: str, placeholder_data: list[dict], output_path: str, slides_to_remove: list[int] = None) -> str:
    prs = Presentation(file_path)
    text_values, image_values, temp_files = prepare_data(placeholder_data)

    for slide in prs.slides:
        process_all_shapes(slide.shapes, slide, text_values, image_values)

    if slides_to_remove:
        slides_to_remove_sorted = sorted(set(slides_to_remove), reverse=True)
        for slide_index in slides_to_remove_sorted:
            if 0 <= slide_index < len(prs.slides):
                slide_id = prs.slides._sldIdLst[slide_index]
                prs.part.drop_rel(slide_id.rId)
                prs.slides._sldIdLst.remove(slide_id)

    prs.save(output_path)
    cleanup_temp_files(temp_files)

def main():
    if len(sys.argv) < 3:
        sys.stderr.write("Usage: python generate.py <file_path> <json_data> [slides_to_remove_json]\n")
        sys.exit(1)
    
    file_path = sys.argv[1]
    json_data = sys.argv[2]
    output_path = sys.argv[3]
    slides_to_remove_json = sys.argv[4] if len(sys.argv) > 3 else '[]'

    try:
        placeholder_data = json.loads(json_data)
        slides_to_remove = json.loads(slides_to_remove_json) if slides_to_remove_json else []
        generate(file_path, placeholder_data, output_path, slides_to_remove)
    except Exception as e:
        sys.stderr.write(f"Error: {e}\n")
        sys.exit(1)

if __name__ == "__main__":
    main()
