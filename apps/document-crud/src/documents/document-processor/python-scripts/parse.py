import sys
import json
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor

def emu_to_pixels(emu_value: int) -> float:
    return round(Emu(emu_value).pt * 96 / 72, 2)

def get_run_color(run) -> RGBColor | None:
    try:
        if run.font.color.rgb:
            return run.font.color.rgb
    except:
        pass
    return None

def get_type_from_color(color: RGBColor | None) -> str | None:
    if color is None:
        return None
    
    r, g, b = color[0], color[1], color[2]
    
    if r > 150 and g < 100 and b < 100:
        return 'text'
    if r < 100 and g < 100 and b > 150:
        return 'image'
    if r < 100 and g > 150 and b < 100:
        return 'map'
    
    return None

def is_placeholder_run(run) -> bool:
    font = run.font
    return font.bold and font.italic and font.underline

def parse_text_frame(text_frame, width: float, height: float, placeholders: dict) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not is_placeholder_run(run):
                continue
            
            key = run.text.strip()
            if not key:
                continue
            
            color = get_run_color(run)
            placeholder_type = get_type_from_color(color)
            
            if placeholder_type is None:
                continue
            
            placeholders[key] = {
                "key": key,
                "type": placeholder_type,
                "width": width,
                "height": height
            }

def parse_shape(shape, placeholders: dict) -> None:
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        if hasattr(shape, 'shapes'):
            for child_shape in shape.shapes:
                parse_shape(child_shape, placeholders)
        return
    
    if hasattr(shape, 'has_table') and shape.has_table:
        for cell in shape.table.iter_cells():
            if hasattr(cell, 'text_frame'):
                width = emu_to_pixels(cell.span_width)
                height = emu_to_pixels(cell.span_height)
                parse_text_frame(cell.text_frame, width, height, placeholders)
        return
    
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        width = emu_to_pixels(shape.width)
        height = emu_to_pixels(shape.height)
        parse_text_frame(shape.text_frame, width, height, placeholders)

def parse_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    placeholders: dict[str, dict] = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            parse_shape(shape, placeholders)

    return list(placeholders.values())

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "No file path provided"}))
        sys.exit(1)

    file_path = sys.argv[1]

    try:
        result = parse_pptx(file_path)
        print(json.dumps(result))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()
