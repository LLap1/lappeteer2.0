import sys
import json
import re
import base64
import tempfile
from io import BytesIO
from pptx import Presentation
from lxml import etree

PLACEHOLDER_PATTERN = re.compile(r'\{\{([^:}]+):([^}:]+)\}\}')
SOFT_BREAK = '\v'
NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

def decode_data_url(data_url: str) -> BytesIO:
    base64_data = data_url.split(',', 1)[1]
    return BytesIO(base64.b64decode(base64_data))

def prepare_data(placeholder_data: list[dict]) -> tuple[dict[str, str], dict[str, list[BytesIO]]]:
    text_values: dict[str, str] = {}
    map_values: dict[str, list[BytesIO]] = {}

    for item in placeholder_data:
        key = item['key'].strip()
        placeholder_type = item['type']
        value = item['value']

        if placeholder_type == 'text':
            text_values[key] = value
        elif placeholder_type == 'map':
            map_values[key] = [decode_data_url(url) for url in json.loads(value)]

    return text_values, map_values

def replace_placeholders_in_text(text: str, text_values: dict[str, str]) -> str:
    def replacer(match):
        key = match.group(1).strip()
        value = text_values.get(key)
        if value is None:
            return match.group(0)
        return value

    return PLACEHOLDER_PATTERN.sub(replacer, text)

def get_paragraph_text(paragraph) -> str:
    text_parts = []
    p_elem = paragraph._p
    
    for child in p_elem:
        tag = etree.QName(child).localname
        if tag == 'r':
            for t_elem in child.findall('.//a:t', NS):
                if t_elem.text:
                    text_parts.append(t_elem.text)
        elif tag == 'br':
            text_parts.append(SOFT_BREAK)
    
    return ''.join(text_parts)

def set_paragraph_text(paragraph, new_text: str) -> None:
    runs = list(paragraph.runs)
    if not runs:
        return
    
    p_elem = paragraph._p
    first_run = runs[0]
    first_run_elem = first_run._r
    
    first_run_copy = etree.fromstring(etree.tostring(first_run_elem))
    for t_elem in first_run_copy.findall('.//a:t', NS):
        first_run_copy.remove(t_elem)
    
    children_to_remove = [child for child in p_elem if etree.QName(child).localname in ('r', 'br')]
    insert_index = list(p_elem).index(children_to_remove[0]) if children_to_remove else len(list(p_elem))
    
    for child in children_to_remove:
        p_elem.remove(child)
    
    parts = new_text.split(SOFT_BREAK)
    
    for i, part in enumerate(parts):
        if i > 0:
            br_elem = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}br')
            p_elem.insert(insert_index, br_elem)
            insert_index += 1
        
        run_elem = etree.fromstring(etree.tostring(first_run_copy))
        t_elem = etree.SubElement(run_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
        t_elem.text = part
        p_elem.insert(insert_index, run_elem)
        insert_index += 1

def replace_text_in_paragraph(paragraph, text_values: dict[str, str]) -> None:
    runs = list(paragraph.runs)
    if not runs:
        return

    full_text = get_paragraph_text(paragraph)
    new_text = replace_placeholders_in_text(full_text, text_values)

    if new_text == full_text:
        return

    set_paragraph_text(paragraph, new_text)

def replace_text_in_frame(text_frame, text_values: dict[str, str]) -> None:
    for paragraph in text_frame.paragraphs:
        replace_text_in_paragraph(paragraph, text_values)

def find_map_placeholder(text: str, map_values: dict[str, list[BytesIO]]) -> str | None:
    matches = PLACEHOLDER_PATTERN.findall(text)
    
    for key, _ in matches:
        key = key.strip()
        if key in map_values:
            return key
    
    return None

def process_shape(shape, slide, text_values: dict[str, str], map_values: dict[str, list[BytesIO]]) -> bool:
    if not hasattr(shape, 'text_frame'):
        return False

    text_frame = shape.text_frame
    map_key = find_map_placeholder(text_frame.text, map_values)

    if map_key:
        images = map_values[map_key]
        left, top, width, height = shape.left, shape.top, shape.width, shape.height

        shape.element.getparent().remove(shape.element)

        for image_buffer in images:
            image_buffer.seek(0)
            slide.shapes.add_picture(image_buffer, left, top, width, height)

        return True

    replace_text_in_frame(text_frame, text_values)
    return False

def process_table(table, text_values: dict[str, str]) -> None:
    for row in table.rows:
        for cell in row.cells:
            if hasattr(cell, 'text_frame'):
                replace_text_in_frame(cell.text_frame, text_values)

def generate(file_path: str, placeholder_data: list[dict]) -> str:
    prs = Presentation(file_path)
    text_values, map_values = prepare_data(placeholder_data)

    for slide in prs.slides:
        shapes_to_process = list(slide.shapes)

        for shape in shapes_to_process:
            if hasattr(shape, 'table'):
                process_table(shape.table, text_values)
            else:
                process_shape(shape, slide, text_values, map_values)

    with tempfile.NamedTemporaryFile(mode='wb', delete=False, suffix='.pptx') as tmp:
        prs.save(tmp)
        return tmp.name

def main():
    if len(sys.argv) < 3:
        sys.stderr.write("Usage: python generate.py <file_path> <json_data>\n")
        sys.exit(1)

    file_path = sys.argv[1]
    json_data = sys.argv[2]

    try:
        placeholder_data = json.loads(json_data)
        output_path = generate(file_path, placeholder_data)
        print(output_path)
    except Exception as e:
        sys.stderr.write(f"Error: {e}\n")
        sys.exit(1)

if __name__ == "__main__":
    main()
