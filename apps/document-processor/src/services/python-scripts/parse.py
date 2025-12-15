"""
Script to extract placeholder values from PowerPoint presentations
Finds all strings in format {{ keyname : type }}
"""
import json
import re
import sys
import traceback
from io import BytesIO
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame
from pptx.util import Length

EMU_PER_INCH = 914400
DPI = 96
PLACEHOLDER_PATTERN = re.compile(r'\{\{\s*([^:}]+?)\s*\:\s*([^}]*?)\s*\}\}')


def emu_to_px(value: Length | int) -> float:
    return round((float(value) / EMU_PER_INCH) * DPI, 2)


def record_placeholder(
    placeholders: Dict[str, Dict[str, str | float]],
    key: str,
    type: str,
    width_emu: Optional[Length | int],
    height_emu: Optional[Length | int]
) -> None:
    if width_emu is None or height_emu is None:
        return
    
    width_px: float = emu_to_px(width_emu)
    height_px: float = emu_to_px(height_emu)
    
    if key not in placeholders:
        placeholders[key] = {
            'key': key,
            'type': type,
            'width': width_px,
            'height': height_px,
        }
        print(f"  Recorded placeholder: {key} - {width_px}x{height_px}px", file=sys.stderr)
    else:
        print(f"  Skipping duplicate placeholder: {key} (already recorded)", file=sys.stderr)


def find_placeholders_in_text_frame(text_frame: TextFrame) -> List[tuple[str, str]]:
    found_placeholders_set: set[tuple[str, str]] = set()
    text = text_frame.text
    
    clean_text = ''.join(
        char for char in text 
        if not (0x200B <= ord(char) <= 0x200F or 0x202A <= ord(char) <= 0x202E or 0x2066 <= ord(char) <= 0x2069)
    )
    
    matches = PLACEHOLDER_PATTERN.findall(clean_text)
    if matches:
        print(f"  Found text: {text}", file=sys.stderr)
        print(f"  Clean text: {clean_text}", file=sys.stderr)
        print(f"  Matches: {matches}", file=sys.stderr)

    for key, type in matches:
        key = key.strip()
        type = type.strip()
        
        if not type and ':' in key:
            parts = key.split(':', 1)
            if len(parts) == 2:
                key = parts[0].strip()
                type = parts[1].strip()
        
        found_placeholders_set.add((key, type))

    return list(found_placeholders_set)


def find_placeholders_in_shape(shape: BaseShape) -> List[tuple[str, str]]:
    if not hasattr(shape, 'text_frame') or not shape.text_frame:
        return []
    
    return find_placeholders_in_text_frame(shape.text_frame)


def process_shape_for_placeholders(
    shape: BaseShape,
    placeholders: Dict[str, Dict[str, str | float]]
) -> None:
    found_placeholders: List[tuple[str, str]] = find_placeholders_in_shape(shape)
    
    if not found_placeholders:
        return
    
    width: Optional[Length | int] = getattr(shape, 'width', None)
    height: Optional[Length | int] = getattr(shape, 'height', None)
    
    for key, type in found_placeholders:
        record_placeholder(placeholders, key, type, width, height)


def process_table_for_placeholders(
    shape: BaseShape,
    placeholders: Dict[str, Dict[str, str | float]]
) -> None:
    if not hasattr(shape, 'table') or not shape.table:
        return
    
    for row in shape.table.rows:
        for cell in row.cells:
            if cell.text_frame:
                found_placeholders = find_placeholders_in_text_frame(cell.text_frame)
                
                if found_placeholders:
                    cell_width = getattr(cell, 'width', None)
                    cell_height = getattr(cell, 'height', None)
                    
                    for key, type in found_placeholders:
                        record_placeholder(placeholders, key, type, cell_width, cell_height)


def extract_placeholders_from_presentation(pptx_data: BytesIO) -> List[Dict[str, str | float]]:
    prs: Presentation = Presentation(pptx_data)
    placeholders: Dict[str, Dict[str, str | float]] = {}
    
    print(f"Analyzing presentation with {len(prs.slides)} slides", file=sys.stderr)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_idx}...", file=sys.stderr)
        
        for shape in slide.shapes:
            if hasattr(shape, 'table') and shape.table:
                process_table_for_placeholders(shape, placeholders)
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                process_shape_for_placeholders(shape, placeholders)
    
    print(f"\nTotal unique placeholders found: {len(placeholders)}", file=sys.stderr)
    
    return sorted(placeholders.values(), key=lambda item: item['key'])


def main() -> None:
    if len(sys.argv) < 2:
        sys.stderr.write("Usage: python3 parse.py <filepath>\n")
        sys.exit(1)
    
    filepath: str = sys.argv[1]
    
    try:
        with open(filepath, 'rb') as f:
            pptx_data: BytesIO = BytesIO(f.read())
    except Exception:
        trace = traceback.format_exc()
        sys.stderr.write(f"Error reading file {filepath}: {trace}\n")
        sys.exit(1)
    
    try:
        placeholders = extract_placeholders_from_presentation(pptx_data)
        output_data = json.dumps(placeholders)
        print(output_data, file=sys.stdout)
        sys.exit(0)
    except Exception:
        trace = traceback.format_exc()
        sys.stderr.write(trace)
        sys.exit(1)


if __name__ == "__main__":
    main()
